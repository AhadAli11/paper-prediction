from docx import Document
from pptx import Presentation
from typing import List, Dict
import re

def extract_topic_blocks(file_path: str) -> List[Dict[str, str]]:
    """
    Extract 'Heading 1' styled text (for .docx) or slide titles (for .pptx) and subsequent text.
    Returns a list of dictionaries with heading as key and following text as value.
    """
    try:
        # Ensure the file is a .docx or .pptx
        if not file_path.lower().endswith(('.docx', '.pptx')):
            raise ValueError("File must be a .docx or .pptx")

        heading_blocks = []

        if file_path.lower().endswith('.docx'):
            # Load the Word document
            doc = Document(file_path)
            current_heading = None
            current_content = []

            # Iterate through paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    if para.style.name == 'Heading 1':
                        # Save previous heading and content if they exist
                        if current_heading and current_content:
                            heading_blocks.append({
                                'heading': current_heading,
                                'content': '\n'.join(current_content).strip()
                            })
                        # Start new heading
                        current_heading = para.text.strip()
                        current_content = []
                    else:
                        # Add to content of current heading
                        current_content.append(para.text.strip())

            # Save the last heading and content if they exist
            if current_heading and current_content:
                heading_blocks.append({
                    'heading': current_heading,
                    'content': '\n'.join(current_content).strip()
                })

        elif file_path.lower().endswith('.pptx'):
            # Load the PowerPoint presentation
            prs = Presentation(file_path)

            # Iterate through slides
            for slide in prs.slides:
                heading = None
                content = []

                # Try to get the slide title from the title placeholder
                if slide.shapes.title and slide.shapes.title.text.strip():
                    heading = slide.shapes.title.text.strip()

                # Extract text from all other shapes in the slide
                for shape in slide.shapes:
                    # Skip the title shape
                    if shape == slide.shapes.title or not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame
                    for para in text_frame.paragraphs:
                        if para.text.strip():
                            content.append(para.text.strip())

                # Add to heading_blocks if a heading was found
                if heading:
                    heading_blocks.append({
                        'heading': heading,
                        'content': '\n'.join(content).strip() if content else ''
                    })

        return heading_blocks

    except Exception as e:
        raise Exception(f"Error extracting headings from {file_path}: {str(e)}")

def extract_questions_from_docx(file_path: str) -> List[str]:
    """
    Extract individual questions from a Word (.docx) quiz or midterm file.
    Supports labels like Q1:, Q1., Question 1:, Question 1., 1., 1:, 1).
    """
    try:
        if not file_path.lower().endswith('.docx'):
            raise ValueError("Only .docx files are supported.")

        doc = Document(file_path)
        full_text = "\n".join(para.text.strip() for para in doc.paragraphs if para.text.strip())

        # Enhanced regex pattern:
        # - Q1:, Q1., Q1), Question 1:, Question 1., Question 1)
        # - 1., 2., 3., 1:, 2:, 3:, 1), 2), 3)
        pattern = r'((?:Q(?:uestion)?\s*\d+|[1-9]\d*)(?:[:.)]))'

        # Split the text into sections
        questions = re.split(pattern, full_text, flags=re.IGNORECASE)

        result = []
        for i in range(1, len(questions), 2):
            label = questions[i].strip()
            content = questions[i + 1].strip() if i + 1 < len(questions) else ""
            result.append(f"{label} {content}")

        return result

    except Exception as e:
        print(f"Error extracting questions from {file_path}: {e}")
        return []

    except Exception as e:
        raise Exception(f"Error extracting questions from {file_path}: {str(e)}")